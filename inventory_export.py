import streamlit as st
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import io
import base64
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import tempfile
import os
from db_operations import get_dashboard_data, get_all_orders, get_all_buyers
from database import get_db_session, Order, Style, Material, ProductionLine, LineAllocation, ProductionEntry

def show_inventory_export():
    """Display the inventory insights export interface"""
    st.title("📊 Inventory Insights Export")
    
    with st.container():
        st.markdown("""
        <div style='background-color: #f0f8ff; padding: 15px; border-radius: 5px; margin-bottom: 20px;'>
            <h3 style='color: #1E3A8A;'>One-click Export to Presentation-Ready Slides</h3>
            <p>Generate professional PowerPoint presentations with key inventory insights, 
            analytics, and visualizations for executive review and decision-making.</p>
        </div>
        """, unsafe_allow_html=True)
    
    # Create tabs for different export options
    tab1, tab2, tab3 = st.tabs(["📊 Inventory Overview", "🏬 Store Analysis", "🔄 Supply Chain"])
    
    with tab1:
        show_inventory_overview_export()
    
    with tab2:
        show_store_analysis_export()
        
    with tab3:
        show_supply_chain_export()

def show_inventory_overview_export():
    """Display the inventory overview export interface"""
    st.subheader("Inventory Overview Export")
    
    # Get inventory data from database
    session = get_db_session()
    
    # Simulate inventory data for different locations
    st.markdown("### Select Data to Include")
    
    # Options for what to include in the export
    include_summary = st.checkbox("Include Inventory Summary", value=True)
    include_trending = st.checkbox("Include Trending Products", value=True)
    include_alerts = st.checkbox("Include Inventory Alerts", value=True)
    include_forecasts = st.checkbox("Include Sales Forecasts", value=True)
    
    # Date range selection
    st.markdown("### Select Date Range")
    col1, col2 = st.columns(2)
    with col1:
        start_date = st.date_input("Start Date", value=pd.to_datetime("2023-01-01"))
    with col2:
        end_date = st.date_input("End Date", value=pd.to_datetime("2023-12-31"))
    
    # Presentation customization
    st.markdown("### Presentation Options")
    company_name = st.text_input("Company Name", value="Voi Jeans Retail India")
    presentation_title = st.text_input("Presentation Title", value="Inventory Insights Report")
    
    theme_options = ["Professional Blue", "Executive Dark", "Retail Vibrant", "Manufacturing Focused", "Supply Chain Themed"]
    selected_theme = st.selectbox("Presentation Theme", theme_options)
    
    include_company_logo = st.checkbox("Include Company Logo", value=True)
    include_date_on_slides = st.checkbox("Include Date on Slides", value=True)
    
    # Preview section
    st.markdown("### Export Preview")
    
    # Generate some sample inventory data for preview
    preview_data = generate_sample_inventory_data()
    
    # Display preview of inventory summary
    if include_summary:
        st.markdown("#### Inventory Summary")
        st.dataframe(preview_data["summary"])
    
    # Display preview of trending products
    if include_trending:
        st.markdown("#### Trending Products")
        st.dataframe(preview_data["trending"])
    
    # Display preview of inventory alerts
    if include_alerts:
        st.markdown("#### Inventory Alerts")
        st.dataframe(preview_data["alerts"])
    
    # Display preview of sales forecasts
    if include_forecasts:
        st.markdown("#### Sales Forecasts")
        st.dataframe(preview_data["forecasts"])
    
    # Export button
    if st.button("🖼️ Export to PowerPoint", type="primary", use_container_width=True):
        # Create PowerPoint presentation
        pptx_file = create_inventory_pptx(
            company_name=company_name, 
            title=presentation_title,
            theme=selected_theme,
            include_logo=include_company_logo,
            include_date=include_date_on_slides,
            data=preview_data,
            include_summary=include_summary,
            include_trending=include_trending,
            include_alerts=include_alerts,
            include_forecasts=include_forecasts
        )
        
        # Create download button
        st.download_button(
            label="📥 Download Presentation",
            data=pptx_file,
            file_name=f"Inventory_Insights_{datetime.now().strftime('%Y%m%d')}.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            use_container_width=True
        )
        
        st.success("Presentation created successfully!")

def show_store_analysis_export():
    """Display the store analysis export interface"""
    st.subheader("Store Analysis Export")
    
    # Options for what to include in the export
    st.markdown("### Select Data to Include")
    include_store_performance = st.checkbox("Include Store Performance Metrics", value=True)
    include_sales_comparison = st.checkbox("Include Sales Comparison", value=True)
    include_customer_demographics = st.checkbox("Include Customer Demographics", value=True)
    include_product_performance = st.checkbox("Include Product Performance by Store", value=True)
    
    # Store selection
    st.markdown("### Select Stores")
    all_stores = [
        "All Stores",
        "Mumbai - Phoenix Mall",
        "Delhi - Select Citywalk",
        "Bangalore - Forum Mall",
        "Chennai - Express Avenue",
        "Hyderabad - Inorbit Mall",
        "Pune - Phoenix Marketcity",
        "Kolkata - South City Mall",
        "Ahmedabad - Alpha One Mall",
        "Jaipur - World Trade Park",
        "Lucknow - Phoenix Palassio"
    ]
    selected_stores = st.multiselect("Select Stores to Include", all_stores, default=["All Stores"])
    
    # Date range selection
    st.markdown("### Select Date Range")
    col1, col2 = st.columns(2)
    with col1:
        start_date = st.date_input("Start Date", value=pd.to_datetime("2023-01-01"), key="store_start_date")
    with col2:
        end_date = st.date_input("End Date", value=pd.to_datetime("2023-12-31"), key="store_end_date")
    
    # Presentation customization
    st.markdown("### Presentation Options")
    company_name = st.text_input("Company Name", value="Voi Jeans Retail India", key="store_company_name")
    presentation_title = st.text_input("Presentation Title", value="Store Performance Analysis", key="store_presentation_title")
    
    theme_options = ["Professional Blue", "Executive Dark", "Retail Vibrant", "Manufacturing Focused", "Supply Chain Themed"]
    selected_theme = st.selectbox("Presentation Theme", theme_options, key="store_theme")
    
    # Preview section
    st.markdown("### Export Preview")
    
    # Generate some sample store data for preview
    preview_data = generate_sample_store_data()
    
    # Display preview of store performance
    if include_store_performance:
        st.markdown("#### Store Performance Metrics")
        st.dataframe(preview_data["performance"])
    
    # Display preview of sales comparison
    if include_sales_comparison:
        st.markdown("#### Sales Comparison")
        st.dataframe(preview_data["sales"])
    
    # Export button
    if st.button("🖼️ Export Store Analysis to PowerPoint", type="primary", use_container_width=True):
        # Create PowerPoint presentation
        pptx_file = create_store_pptx(
            company_name=company_name, 
            title=presentation_title,
            theme=selected_theme,
            data=preview_data,
            include_performance=include_store_performance,
            include_sales=include_sales_comparison,
            include_demographics=include_customer_demographics,
            include_product=include_product_performance
        )
        
        # Create download button
        st.download_button(
            label="📥 Download Store Analysis",
            data=pptx_file,
            file_name=f"Store_Analysis_{datetime.now().strftime('%Y%m%d')}.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            use_container_width=True
        )
        
        st.success("Store Analysis presentation created successfully!")

def show_supply_chain_export():
    """Display the supply chain export interface"""
    st.subheader("Supply Chain Export")
    
    # Options for what to include in the export
    st.markdown("### Select Data to Include")
    include_supplier_performance = st.checkbox("Include Supplier Performance", value=True)
    include_logistics = st.checkbox("Include Logistics Overview", value=True)
    include_lead_times = st.checkbox("Include Lead Time Analysis", value=True)
    include_inventory_flow = st.checkbox("Include Inventory Flow Visualization", value=True)
    
    # Date range selection
    st.markdown("### Select Date Range")
    col1, col2 = st.columns(2)
    with col1:
        start_date = st.date_input("Start Date", value=pd.to_datetime("2023-01-01"), key="supply_start_date")
    with col2:
        end_date = st.date_input("End Date", value=pd.to_datetime("2023-12-31"), key="supply_end_date")
    
    # Presentation customization
    st.markdown("### Presentation Options")
    company_name = st.text_input("Company Name", value="Voi Jeans Retail India", key="supply_company_name")
    presentation_title = st.text_input("Presentation Title", value="Supply Chain Analysis", key="supply_presentation_title")
    
    theme_options = ["Professional Blue", "Executive Dark", "Retail Vibrant", "Manufacturing Focused", "Supply Chain Themed"]
    selected_theme = st.selectbox("Presentation Theme", theme_options, key="supply_theme")
    
    # Preview section
    st.markdown("### Export Preview")
    
    # Generate some sample supply chain data for preview
    preview_data = generate_sample_supply_chain_data()
    
    # Display preview of supplier performance
    if include_supplier_performance:
        st.markdown("#### Supplier Performance")
        st.dataframe(preview_data["suppliers"])
    
    # Display preview of logistics
    if include_logistics:
        st.markdown("#### Logistics Overview")
        st.dataframe(preview_data["logistics"])
    
    # Export button
    if st.button("🖼️ Export Supply Chain Analysis to PowerPoint", type="primary", use_container_width=True):
        # Create PowerPoint presentation
        pptx_file = create_supply_chain_pptx(
            company_name=company_name, 
            title=presentation_title,
            theme=selected_theme,
            data=preview_data,
            include_suppliers=include_supplier_performance,
            include_logistics=include_logistics,
            include_lead_times=include_lead_times,
            include_flow=include_inventory_flow
        )
        
        # Create download button
        st.download_button(
            label="📥 Download Supply Chain Analysis",
            data=pptx_file,
            file_name=f"Supply_Chain_Analysis_{datetime.now().strftime('%Y%m%d')}.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            use_container_width=True
        )
        
        st.success("Supply Chain Analysis presentation created successfully!")

def generate_sample_inventory_data():
    """Generate sample inventory data for preview and export"""
    # Inventory summary data
    summary_data = {
        'Category': ['Denim Jeans', 'T-Shirts', 'Shirts', 'Jackets', 'Accessories'],
        'Total Units': [1250, 980, 750, 320, 1560],
        'Available Units': [950, 780, 650, 280, 1360],
        'Reserved Units': [300, 200, 100, 40, 200],
        'Avg. Selling Price': ['₹2,499', '₹999', '₹1,499', '₹3,499', '₹599'],
        'Total Value': ['₹23,74,050', '₹7,76,220', '₹9,73,500', '₹9,79,200', '₹8,14,840']
    }
    
    # Trending products data
    trending_data = {
        'Product': ['Slim Fit Blue Jeans', 'Black Distressed Denim', 'White Logo T-Shirt', 'Denim Jacket', 'Canvas Belt'],
        'Sales Last 7 Days': [89, 76, 62, 43, 41],
        'Sales Last 30 Days': [356, 298, 241, 182, 175],
        'Growth Rate': ['+12.4%', '+8.7%', '+5.2%', '+15.1%', '+3.8%'],
        'Stock Level': ['Medium', 'Low', 'Medium', 'Medium', 'High'],
        'Reorder Status': ['Not Required', 'Reorder', 'Not Required', 'Not Required', 'Not Required']
    }
    
    # Inventory alerts data
    alerts_data = {
        'Product': ['Black Distressed Denim', 'Slim Fit Grey Jeans', 'Basic White T-Shirt', 'Leather Belt', 'Blue Casual Shirt'],
        'Current Stock': [28, 22, 35, 12, 19],
        'Reorder Point': [30, 25, 40, 15, 20],
        'Status': ['Below Reorder Point', 'Below Reorder Point', 'Below Reorder Point', 'Below Reorder Point', 'Below Reorder Point'],
        'Suggested Order': [50, 50, 100, 30, 40],
        'Expected Delivery': ['2023-07-15', '2023-07-18', '2023-07-12', '2023-07-20', '2023-07-14']
    }
    
    # Sales forecasts data
    forecasts_data = {
        'Category': ['Denim Jeans', 'T-Shirts', 'Shirts', 'Jackets', 'Accessories'],
        'Current Month Forecast': [350, 280, 220, 90, 420],
        'Next Month Forecast': [380, 320, 240, 120, 450],
        'Growth': ['+8.6%', '+14.3%', '+9.1%', '+33.3%', '+7.1%'],
        'Confidence Level': ['High', 'Medium', 'Medium', 'Low', 'High'],
        'Stock Coverage': ['3.2 months', '2.8 months', '3.0 months', '3.1 months', '3.2 months']
    }
    
    return {
        'summary': pd.DataFrame(summary_data),
        'trending': pd.DataFrame(trending_data),
        'alerts': pd.DataFrame(alerts_data),
        'forecasts': pd.DataFrame(forecasts_data)
    }

def generate_sample_store_data():
    """Generate sample store data for preview and export"""
    # Store performance data
    performance_data = {
        'Store': ['Mumbai - Phoenix Mall', 'Delhi - Select Citywalk', 'Bangalore - Forum Mall', 'Chennai - Express Avenue', 'Hyderabad - Inorbit Mall'],
        'Monthly Sales': ['₹42,56,780', '₹38,92,450', '₹35,67,890', '₹29,45,670', '₹31,28,940'],
        'Sales Growth': ['+8.4%', '+5.2%', '+7.6%', '+3.8%', '+6.2%'],
        'Footfall': [12500, 11800, 10500, 9200, 9800],
        'Conversion Rate': ['12.4%', '11.8%', '13.2%', '10.9%', '11.5%'],
        'Avg. Basket Size': ['₹2,750', '₹2,890', '₹2,650', '₹2,430', '₹2,520']
    }
    
    # Sales comparison data
    sales_data = {
        'Store': ['Mumbai - Phoenix Mall', 'Delhi - Select Citywalk', 'Bangalore - Forum Mall', 'Chennai - Express Avenue', 'Hyderabad - Inorbit Mall'],
        'Denim Sales': ['₹18,45,670', '₹16,78,900', '₹15,34,560', '₹12,67,890', '₹13,56,780'],
        'T-Shirt Sales': ['₹8,92,340', '₹7,89,650', '₹8,12,450', '₹6,78,900', '₹7,23,450'],
        'Shirt Sales': ['₹7,45,670', '₹6,89,230', '₹6,45,670', '₹5,67,890', '₹5,98,760'],
        'Jacket Sales': ['₹5,23,450', '₹4,78,900', '₹3,98,760', '₹2,87,650', '₹3,12,340'],
        'Accessories Sales': ['₹2,49,650', '₹2,55,770', '₹1,76,450', '₹1,43,340', '₹1,37,610']
    }
    
    return {
        'performance': pd.DataFrame(performance_data),
        'sales': pd.DataFrame(sales_data)
    }

def generate_sample_supply_chain_data():
    """Generate sample supply chain data for preview and export"""
    # Supplier performance data
    supplier_data = {
        'Supplier': ['Raymond Textiles', 'Arvind Mills', 'JCT Limited', 'Vardhman Textiles', 'Bombay Dyeing'],
        'On-Time Delivery': ['92%', '88%', '94%', '90%', '86%'],
        'Quality Rating': ['4.5/5', '4.2/5', '4.7/5', '4.3/5', '4.0/5'],
        'Cost Efficiency': ['High', 'Medium', 'High', 'Medium', 'Medium'],
        'Lead Time (Days)': [12, 15, 10, 14, 18],
        'Last Order Date': ['2023-06-15', '2023-06-10', '2023-06-18', '2023-06-12', '2023-06-08']
    }
    
    # Logistics data
    logistics_data = {
        'Route': ['Mumbai to Delhi', 'Mumbai to Bangalore', 'Delhi to Chennai', 'Bangalore to Hyderabad', 'Chennai to Mumbai'],
        'Transit Time (Days)': [2, 3, 4, 2, 3],
        'Avg. Cost per Shipment': ['₹45,000', '₹62,000', '₹78,000', '₹38,000', '₹58,000'],
        'Reliability Score': ['4.6/5', '4.3/5', '4.2/5', '4.7/5', '4.4/5'],
        'Carrier': ['Blue Dart', 'DHL', 'FedEx', 'Blue Dart', 'DHL'],
        'Last Shipment Date': ['2023-06-20', '2023-06-18', '2023-06-15', '2023-06-19', '2023-06-17']
    }
    
    return {
        'suppliers': pd.DataFrame(supplier_data),
        'logistics': pd.DataFrame(logistics_data)
    }

def create_inventory_pptx(company_name, title, theme, include_logo, include_date, data, 
                         include_summary, include_trending, include_alerts, include_forecasts):
    """Create a PowerPoint presentation with inventory insights"""
    # Create a presentation
    prs = Presentation()
    
    # Get theme colors
    theme_colors = get_theme_colors(theme)
    
    # Create the title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    
    # Add title
    title = slide.shapes.title
    title.text = title
    subtitle = slide.placeholders[1]
    subtitle.text = f"{company_name}\n{datetime.now().strftime('%B %d, %Y')}"
    
    # Add logo if requested
    if include_logo:
        # In a real implementation, would load the company logo
        # Here we'll just add a text box as a placeholder
        left = Inches(9)
        top = Inches(0.5)
        width = Inches(1)
        height = Inches(1)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = "LOGO"
    
    # Add inventory summary slide if requested
    if include_summary:
        add_inventory_summary_slide(prs, data["summary"], theme_colors, include_date, company_name)
    
    # Add trending products slide if requested
    if include_trending:
        add_trending_products_slide(prs, data["trending"], theme_colors, include_date, company_name)
    
    # Add inventory alerts slide if requested
    if include_alerts:
        add_inventory_alerts_slide(prs, data["alerts"], theme_colors, include_date, company_name)
    
    # Add sales forecasts slide if requested
    if include_forecasts:
        add_sales_forecasts_slide(prs, data["forecasts"], theme_colors, include_date, company_name)
    
    # Save the presentation to a BytesIO object
    pptx_file = io.BytesIO()
    prs.save(pptx_file)
    pptx_file.seek(0)
    
    return pptx_file

def create_store_pptx(company_name, title, theme, data, 
                     include_performance, include_sales, include_demographics, include_product):
    """Create a PowerPoint presentation with store analysis"""
    # Create a presentation
    prs = Presentation()
    
    # Get theme colors
    theme_colors = get_theme_colors(theme)
    
    # Create the title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    
    # Add title
    title = slide.shapes.title
    title.text = title
    subtitle = slide.placeholders[1]
    subtitle.text = f"{company_name}\n{datetime.now().strftime('%B %d, %Y')}"
    
    # Add store performance slide if requested
    if include_performance:
        add_store_performance_slide(prs, data["performance"], theme_colors, True, company_name)
    
    # Add sales comparison slide if requested
    if include_sales:
        add_sales_comparison_slide(prs, data["sales"], theme_colors, True, company_name)
    
    # Save the presentation to a BytesIO object
    pptx_file = io.BytesIO()
    prs.save(pptx_file)
    pptx_file.seek(0)
    
    return pptx_file

def create_supply_chain_pptx(company_name, title, theme, data, 
                            include_suppliers, include_logistics, include_lead_times, include_flow):
    """Create a PowerPoint presentation with supply chain analysis"""
    # Create a presentation
    prs = Presentation()
    
    # Get theme colors
    theme_colors = get_theme_colors(theme)
    
    # Create the title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    
    # Add title
    title = slide.shapes.title
    title.text = title
    subtitle = slide.placeholders[1]
    subtitle.text = f"{company_name}\n{datetime.now().strftime('%B %d, %Y')}"
    
    # Add supplier performance slide if requested
    if include_suppliers:
        add_supplier_performance_slide(prs, data["suppliers"], theme_colors, True, company_name)
    
    # Add logistics slide if requested
    if include_logistics:
        add_logistics_slide(prs, data["logistics"], theme_colors, True, company_name)
    
    # Save the presentation to a BytesIO object
    pptx_file = io.BytesIO()
    prs.save(pptx_file)
    pptx_file.seek(0)
    
    return pptx_file

def get_theme_colors(theme):
    """Get the colors for a specific theme"""
    if theme == "Professional Blue":
        return {
            "primary": RGBColor(31, 73, 125),
            "secondary": RGBColor(79, 129, 189),
            "text": RGBColor(0, 0, 0),
            "background": RGBColor(255, 255, 255),
            "accent1": RGBColor(192, 80, 77),
            "accent2": RGBColor(155, 187, 89)
        }
    elif theme == "Executive Dark":
        return {
            "primary": RGBColor(68, 84, 106),
            "secondary": RGBColor(68, 114, 196),
            "text": RGBColor(255, 255, 255),
            "background": RGBColor(35, 35, 35),
            "accent1": RGBColor(237, 125, 49),
            "accent2": RGBColor(165, 165, 165)
        }
    elif theme == "Retail Vibrant":
        return {
            "primary": RGBColor(91, 155, 213),
            "secondary": RGBColor(237, 125, 49),
            "text": RGBColor(0, 0, 0),
            "background": RGBColor(248, 248, 248),
            "accent1": RGBColor(255, 192, 0),
            "accent2": RGBColor(112, 173, 71)
        }
    elif theme == "Manufacturing Focused":
        return {
            "primary": RGBColor(68, 114, 196),
            "secondary": RGBColor(112, 173, 71),
            "text": RGBColor(0, 0, 0),
            "background": RGBColor(242, 242, 242),
            "accent1": RGBColor(255, 192, 0),
            "accent2": RGBColor(237, 125, 49)
        }
    elif theme == "Supply Chain Themed":
        return {
            "primary": RGBColor(0, 112, 192),
            "secondary": RGBColor(0, 176, 80),
            "text": RGBColor(0, 0, 0),
            "background": RGBColor(255, 255, 255),
            "accent1": RGBColor(255, 192, 0),
            "accent2": RGBColor(192, 0, 0)
        }
    else:
        # Default theme
        return {
            "primary": RGBColor(31, 73, 125),
            "secondary": RGBColor(79, 129, 189),
            "text": RGBColor(0, 0, 0),
            "background": RGBColor(255, 255, 255),
            "accent1": RGBColor(192, 80, 77),
            "accent2": RGBColor(155, 187, 89)
        }

def add_inventory_summary_slide(prs, df, theme_colors, include_date, company_name):
    """Add a slide with inventory summary data"""
    # Create the slide
    slide_layout = prs.slide_layouts[5]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title = slide.shapes.title
    title.text = "Inventory Summary"
    
    # Add date and company name if requested
    if include_date:
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(4)
        height = Inches(0.5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = f"{company_name} - {datetime.now().strftime('%B %d, %Y')}"
    
    # Convert DataFrame to a table
    x, y, cx, cy = Inches(0.5), Inches(1.5), Inches(9), Inches(5)
    shape = slide.shapes.add_table(
        rows=len(df) + 1,
        cols=len(df.columns),
        left=x, top=y, width=cx, height=cy
    )
    table = shape.table
    
    # Add headers
    for i, column in enumerate(df.columns):
        cell = table.cell(0, i)
        cell.text = column
        
        # Style the header cells
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.bold = True
                run.font.size = Pt(12)
    
    # Add data
    for row in range(len(df)):
        for col in range(len(df.columns)):
            cell = table.cell(row + 1, col)
            cell.text = str(df.iloc[row, col])
            
            # Style the data cells
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(10)
    
    return slide

def add_trending_products_slide(prs, df, theme_colors, include_date, company_name):
    """Add a slide with trending products data"""
    # Create the slide
    slide_layout = prs.slide_layouts[5]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title = slide.shapes.title
    title.text = "Trending Products"
    
    # Add date and company name if requested
    if include_date:
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(4)
        height = Inches(0.5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = f"{company_name} - {datetime.now().strftime('%B %d, %Y')}"
    
    # Convert DataFrame to a table
    x, y, cx, cy = Inches(0.5), Inches(1.5), Inches(9), Inches(5)
    shape = slide.shapes.add_table(
        rows=len(df) + 1,
        cols=len(df.columns),
        left=x, top=y, width=cx, height=cy
    )
    table = shape.table
    
    # Add headers
    for i, column in enumerate(df.columns):
        cell = table.cell(0, i)
        cell.text = column
        
        # Style the header cells
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.bold = True
                run.font.size = Pt(12)
    
    # Add data
    for row in range(len(df)):
        for col in range(len(df.columns)):
            cell = table.cell(row + 1, col)
            cell.text = str(df.iloc[row, col])
            
            # Style the data cells
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(10)
    
    return slide

def add_inventory_alerts_slide(prs, df, theme_colors, include_date, company_name):
    """Add a slide with inventory alerts data"""
    # Create the slide
    slide_layout = prs.slide_layouts[5]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title = slide.shapes.title
    title.text = "Inventory Alerts"
    
    # Add date and company name if requested
    if include_date:
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(4)
        height = Inches(0.5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = f"{company_name} - {datetime.now().strftime('%B %d, %Y')}"
    
    # Convert DataFrame to a table
    x, y, cx, cy = Inches(0.5), Inches(1.5), Inches(9), Inches(5)
    shape = slide.shapes.add_table(
        rows=len(df) + 1,
        cols=len(df.columns),
        left=x, top=y, width=cx, height=cy
    )
    table = shape.table
    
    # Add headers
    for i, column in enumerate(df.columns):
        cell = table.cell(0, i)
        cell.text = column
        
        # Style the header cells
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.bold = True
                run.font.size = Pt(12)
    
    # Add data
    for row in range(len(df)):
        for col in range(len(df.columns)):
            cell = table.cell(row + 1, col)
            cell.text = str(df.iloc[row, col])
            
            # Style the data cells
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(10)
    
    return slide

def add_sales_forecasts_slide(prs, df, theme_colors, include_date, company_name):
    """Add a slide with sales forecasts data"""
    # Create the slide
    slide_layout = prs.slide_layouts[5]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title = slide.shapes.title
    title.text = "Sales Forecasts"
    
    # Add date and company name if requested
    if include_date:
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(4)
        height = Inches(0.5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = f"{company_name} - {datetime.now().strftime('%B %d, %Y')}"
    
    # Convert DataFrame to a table
    x, y, cx, cy = Inches(0.5), Inches(1.5), Inches(9), Inches(5)
    shape = slide.shapes.add_table(
        rows=len(df) + 1,
        cols=len(df.columns),
        left=x, top=y, width=cx, height=cy
    )
    table = shape.table
    
    # Add headers
    for i, column in enumerate(df.columns):
        cell = table.cell(0, i)
        cell.text = column
        
        # Style the header cells
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.bold = True
                run.font.size = Pt(12)
    
    # Add data
    for row in range(len(df)):
        for col in range(len(df.columns)):
            cell = table.cell(row + 1, col)
            cell.text = str(df.iloc[row, col])
            
            # Style the data cells
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(10)
    
    return slide

def add_store_performance_slide(prs, df, theme_colors, include_date, company_name):
    """Add a slide with store performance data"""
    # Create the slide
    slide_layout = prs.slide_layouts[5]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title = slide.shapes.title
    title.text = "Store Performance Metrics"
    
    # Add date and company name if requested
    if include_date:
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(4)
        height = Inches(0.5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = f"{company_name} - {datetime.now().strftime('%B %d, %Y')}"
    
    # Convert DataFrame to a table
    x, y, cx, cy = Inches(0.5), Inches(1.5), Inches(9), Inches(5)
    shape = slide.shapes.add_table(
        rows=len(df) + 1,
        cols=len(df.columns),
        left=x, top=y, width=cx, height=cy
    )
    table = shape.table
    
    # Add headers
    for i, column in enumerate(df.columns):
        cell = table.cell(0, i)
        cell.text = column
        
        # Style the header cells
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.bold = True
                run.font.size = Pt(12)
    
    # Add data
    for row in range(len(df)):
        for col in range(len(df.columns)):
            cell = table.cell(row + 1, col)
            cell.text = str(df.iloc[row, col])
            
            # Style the data cells
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(10)
    
    return slide

def add_sales_comparison_slide(prs, df, theme_colors, include_date, company_name):
    """Add a slide with sales comparison data"""
    # Create the slide
    slide_layout = prs.slide_layouts[5]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title = slide.shapes.title
    title.text = "Sales Comparison by Category"
    
    # Add date and company name if requested
    if include_date:
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(4)
        height = Inches(0.5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = f"{company_name} - {datetime.now().strftime('%B %d, %Y')}"
    
    # Convert DataFrame to a table
    x, y, cx, cy = Inches(0.5), Inches(1.5), Inches(9), Inches(5)
    shape = slide.shapes.add_table(
        rows=len(df) + 1,
        cols=len(df.columns),
        left=x, top=y, width=cx, height=cy
    )
    table = shape.table
    
    # Add headers
    for i, column in enumerate(df.columns):
        cell = table.cell(0, i)
        cell.text = column
        
        # Style the header cells
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.bold = True
                run.font.size = Pt(12)
    
    # Add data
    for row in range(len(df)):
        for col in range(len(df.columns)):
            cell = table.cell(row + 1, col)
            cell.text = str(df.iloc[row, col])
            
            # Style the data cells
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(10)
    
    return slide

def add_supplier_performance_slide(prs, df, theme_colors, include_date, company_name):
    """Add a slide with supplier performance data"""
    # Create the slide
    slide_layout = prs.slide_layouts[5]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title = slide.shapes.title
    title.text = "Supplier Performance Metrics"
    
    # Add date and company name if requested
    if include_date:
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(4)
        height = Inches(0.5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = f"{company_name} - {datetime.now().strftime('%B %d, %Y')}"
    
    # Convert DataFrame to a table
    x, y, cx, cy = Inches(0.5), Inches(1.5), Inches(9), Inches(5)
    shape = slide.shapes.add_table(
        rows=len(df) + 1,
        cols=len(df.columns),
        left=x, top=y, width=cx, height=cy
    )
    table = shape.table
    
    # Add headers
    for i, column in enumerate(df.columns):
        cell = table.cell(0, i)
        cell.text = column
        
        # Style the header cells
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.bold = True
                run.font.size = Pt(12)
    
    # Add data
    for row in range(len(df)):
        for col in range(len(df.columns)):
            cell = table.cell(row + 1, col)
            cell.text = str(df.iloc[row, col])
            
            # Style the data cells
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(10)
    
    return slide

def add_logistics_slide(prs, df, theme_colors, include_date, company_name):
    """Add a slide with logistics data"""
    # Create the slide
    slide_layout = prs.slide_layouts[5]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title = slide.shapes.title
    title.text = "Logistics Overview"
    
    # Add date and company name if requested
    if include_date:
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(4)
        height = Inches(0.5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = f"{company_name} - {datetime.now().strftime('%B %d, %Y')}"
    
    # Convert DataFrame to a table
    x, y, cx, cy = Inches(0.5), Inches(1.5), Inches(9), Inches(5)
    shape = slide.shapes.add_table(
        rows=len(df) + 1,
        cols=len(df.columns),
        left=x, top=y, width=cx, height=cy
    )
    table = shape.table
    
    # Add headers
    for i, column in enumerate(df.columns):
        cell = table.cell(0, i)
        cell.text = column
        
        # Style the header cells
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.bold = True
                run.font.size = Pt(12)
    
    # Add data
    for row in range(len(df)):
        for col in range(len(df.columns)):
            cell = table.cell(row + 1, col)
            cell.text = str(df.iloc[row, col])
            
            # Style the data cells
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(10)
    
    return slide